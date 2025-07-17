import streamlit as st
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml import parse_xml
from io import BytesIO
from copy import deepcopy
import tempfile

st.set_page_config(page_title="PowerPoint Content Transfer", layout="centered")
st.title("📊 PowerPoint Content Transfer")
st.markdown("Upload a PowerPoint presentation containing the source content (File A), which will be transferred into a PowerPoint presentation formatted with the company’s official template (File B)")

file_a = st.file_uploader("Upload Source Content PowerPoint (File A)", type=["pptx"])
file_b = st.file_uploader("Upload Company Template PowerPoint (File B)", type=["pptx"])

# === Helper Functions ===
def recursively_ungroup_shapes(slide):
    while True:
        group_shapes = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.GROUP]
        if not group_shapes:
            break

        for group in group_shapes:
            for subshape in group.shapes:
                try:
                    new_el = deepcopy(subshape.element)
                    slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
                except:
                    continue
            try:
                group._element.getparent().remove(group._element)
            except:
                continue

def ungroup_all_shapes(prs):
    for slide in prs.slides:
        recursively_ungroup_shapes(slide)

def copy_shapes_exact(source_slide, target_slide):
    for shape in source_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_stream = BytesIO(shape.image.blob)
            target_slide.shapes.add_picture(
                image_stream, shape.left, shape.top, shape.width, shape.height)
        else:
            try:
                new_el = deepcopy(shape.element)
                target_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
            except:
                continue

def center_shapes_on_slide(slide, slide_width, slide_height):
    shapes = [shape for shape in slide.shapes 
              if shape.shape_type != MSO_SHAPE_TYPE.GROUP and 
                 shape.left is not None and shape.top is not None and
                 shape.width is not None and shape.height is not None]

    if not shapes:
        return

    lefts = [shape.left for shape in shapes]
    tops = [shape.top for shape in shapes]
    rights = [shape.left + shape.width for shape in shapes]
    bottoms = [shape.top + shape.height for shape in shapes]

    min_left = min(lefts)
    min_top = min(tops)
    max_right = max(rights)
    max_bottom = max(bottoms)

    group_width = max_right - min_left
    group_height = max_bottom - min_top

    group_shape_xml = f"""
    <p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
             xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:nvGrpSpPr>
        <p:cNvPr id="1" name="Group 1"/>
        <p:cNvGrpSpPr/>
        <p:nvPr/>
      </p:nvGrpSpPr>
      <p:grpSpPr>
        <a:xfrm>
          <a:off x="{min_left}" y="{min_top}"/>
          <a:ext cx="{group_width}" cy="{group_height}"/>
          <a:chOff x="{min_left}" y="{min_top}"/>
          <a:chExt cx="{group_width}" cy="{group_height}"/>
        </a:xfrm>
      </p:grpSpPr>
    </p:grpSp>
    """
    group_element = parse_xml(group_shape_xml)

    for shape in shapes:
        slide.shapes._spTree.remove(shape.element)
        group_element.append(shape.element)

    slide.shapes._spTree.append(group_element)
    group_shape = slide.shapes[-1]

    new_left = int((slide_width - group_width) / 2)
    new_top = int((slide_height - group_height) / 2)

    xfrm = group_shape.element.find(".//a:xfrm", namespaces={"a": "http://schemas.openxmlformats.org/drawingml/2006/main"})
    xfrm.find("a:off", namespaces={"a": xfrm.nsmap["a"]}).set("x", str(new_left))
    xfrm.find("a:off", namespaces={"a": xfrm.nsmap["a"]}).set("y", str(new_top))

# === Main Processing ===
if file_a and file_b:
    try:
        prs_content = Presentation(file_a)
        prs_template = Presentation(file_b)

        ungroup_all_shapes(prs_content)

        template_layout = prs_template.slides[0].slide_layout
        while len(prs_template.slides) > 0:
            rId = prs_template.slides._sldIdLst[0].rId
            prs_template.part.drop_rel(rId)
            del prs_template.slides._sldIdLst[0]

        for slide in prs_content.slides:
            new_slide = prs_template.slides.add_slide(template_layout)
            copy_shapes_exact(slide, new_slide)

        slide_w = prs_template.slide_width
        slide_h = prs_template.slide_height

        for slide in prs_template.slides:
            center_shapes_on_slide(slide, slide_w, slide_h)

        # Save output
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
            prs_template.save(tmp_file.name)
            tmp_file.seek(0)
            st.success("✅ Processing complete!")
            st.download_button("📥 Download Processed PowerPoint file", tmp_file.read(), file_name="Processed_Presentation.pptx")

    except Exception as e:
        st.error(f"❌ An error occurred: {e}")
